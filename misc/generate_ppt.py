import collections
import collections.abc
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

def create_ppt():
    prs = Presentation()

    # Slide 1: Title Slide (Layout 0)
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Comprehensive NIDS Project Report"
    subtitle.text = "A Hybrid Zero-Day NetGuard Built on Golang and Python AI\n\nAuthor: Gaurav Prasad\nMarch 31, 2026"

    # Slide 2: Abstract / Executive Summary (Layout 1)
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Abstract / Executive Summary"
    tf = slide.placeholders[1].text_frame
    tf.text = "Modern threats bypass traditional signature-matching systems."
    tf.add_paragraph().text = "AI solutions are effective but lack high-speed throughput capability."
    tf.add_paragraph().text = "NetGuard uses a hybrid approach: Golang for wire-speed packet capture & IP blacklisting."
    tf.add_paragraph().text = "Python FastAPI + Scikit-Learn for asynchronous Zero-Day anomaly detection (Isolation Forest)."
    tf.add_paragraph().text = "Telemetry and alerts are streamed via SSE to a live analytics dashboard."

    # Slide 3: Introduction & Problem Statement
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Introduction & Problem Statement"
    tf = slide.placeholders[1].text_frame
    tf.text = "Traditional NIDS rely heavily on static signatures, vulnerable to polymorphic changes."
    tf.add_paragraph().text = "Deep Packet Inspection (DPI) creates massive bottlenecks at 10-40 Gbps."
    tf.add_paragraph().text = "Python's Global Interpreter Lock (GIL) limits pure AI-driven active network interception."
    tf.add_paragraph().text = "Flow-Based Analysis evaluates mathematical shapes of connections instead of payload decoding."

    # Slide 4: Project Objectives
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Project Objectives"
    tf = slide.placeholders[1].text_frame
    tf.text = "1. Develop a line-speed packet parser using low-level network hooks (Golang)."
    tf.add_paragraph().text = "2. Implement high-efficiency deterministic rule-matching for millions of IPs."
    tf.add_paragraph().text = "3. Integrate unsupervised Machine Learning to detect zero-day flow anomalies."
    tf.add_paragraph().text = "4. Deliver a seamless user experience via a real-time web interface (SSE)."

    # Slide 5: Methodology 1/3 - Packet Capture
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Methodology: Packet Capture Engine"
    tf = slide.placeholders[1].text_frame
    tf.text = "Implemented in Golang for memory safety and Goroutine concurrency."
    tf.add_paragraph().text = "Binds to OS network interface using libpcap wrappers (gopacket)."
    tf.add_paragraph().text = "Decodes raw Ethernet frames sequentially into IPv4 and TCP/UDP sub-layers."
    tf.add_paragraph().text = "Concurrent hash-map tracks active bi-directional network states."
    tf.add_paragraph().text = "Summarizes metadata upon connection termination or temporal timeout."

    # Slide 6: Methodology 2/3 - Deterministic Blacklisting
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Methodology: High-Speed Blacklisting"
    tf = slide.placeholders[1].text_frame
    tf.text = "Mass cross-referencing against 4 million indicators within the packet pipeline."
    tf.add_paragraph().text = "Utilizes a custom Trie (Prefix) Tree for memory structuring."
    tf.add_paragraph().text = "IPv4 addresses queried using bitwise edge hops."
    tf.add_paragraph().text = "Operates in O(K) constant time complexity (max 32 edges)."
    tf.add_paragraph().text = "Lookup speeds remain mathematically immune to threat-list size."

    # Slide 7: Methodology 3/3 - AI Engine
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Methodology: Zero-Day AI Engine"
    tf = slide.placeholders[1].text_frame
    tf.text = "Machine Learning operates asynchronously inside isolated Python FastAPI backend."
    tf.add_paragraph().text = "Go engine offloads network features via RESTful HTTP POST mechanism."
    tf.add_paragraph().text = "Hosts pre-trained Scikit-Learn Isolation Forest model."
    tf.add_paragraph().text = "Identifies unexpected traffic profiles without explicit signatures."
    tf.add_paragraph().text = "Yields outlier classification labels (-1) for distinct deviations."

    # Slide 8: Diagram 1 - High-Level Architecture (Placeholder)
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "System Architecture Diagram"
    tf = slide.placeholders[1].text_frame
    tf.text = "[ PLACEHOLDER FOR ARCHITECTURE DIAGRAM ]"
    tf.add_paragraph().text = "Diagram depicting Golang NIDS Engine, Python FastAPI, and Analytics Dashboard."

    # Slide 9: Diagram 2 - Component Interaction (Placeholder)
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Component Interaction Flow"
    tf = slide.placeholders[1].text_frame
    tf.text = "[ PLACEHOLDER FOR INTERACTION DIAGRAM ]"
    tf.add_paragraph().text = "Sequence diagram detailing packet ingress -> rule validation -> flow classification."

    # Slide 10: Results & Findings
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Results & Findings"
    tf = slide.placeholders[1].text_frame
    tf.text = "Memory Footprint: Golang packet tracker maintained < 20MB overhead."
    tf.add_paragraph().text = "Execution Latency: API inference returned predictions within 4-8ms."
    tf.add_paragraph().text = "Trie Tree Scaling: Zero packet-forwarding degradation despite 4-million IP list."
    tf.add_paragraph().text = "Stream Stability: Real-time SSE proved resilient against UI freezing during restarts."

    # Slide 11: Discussion
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Discussion & Analysis"
    tf = slide.placeholders[1].text_frame
    tf.text = "Hybrid microservices model effectively solved structural ML VS Packet analysis conflicts."
    tf.add_paragraph().text = "Overcame asynchronous SSE tracking pointer bugs by utilizing active-reset."
    tf.add_paragraph().text = "Current Limitation: Synchronous HTTP POST introduces high API coupling and backpressure."

    # Slide 12: Conclusion & Recommendations
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Conclusion & Recommendations"
    tf = slide.placeholders[1].text_frame
    tf.text = "Successfully engineered a line-rate capable, intelligent network monitor."
    tf.add_paragraph().text = "Recommendation 1: Transition to High-Availability Message Brokers (Redis, RabbitMQ)."
    tf.add_paragraph().text = "Recommendation 2: Use ElasticSearch or Prometheus for time-series logging."
    tf.add_paragraph().text = "Recommendation 3: Implement active response via eBPF kernel hooks to drop payloads."

    prs.save("Project_Report_Final.pptx")
    print("Successfully generated Project_Report_Final.pptx")

if __name__ == '__main__':
    create_ppt()
